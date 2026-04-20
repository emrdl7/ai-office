"""Tool Registry — Job step에서 사용 가능한 도구 목록과 실행기를 관리한다.

새 도구 추가: tools/ 디렉토리에 <tool_id>.py 파일을 만들거나
              이 파일의 _BUILTIN_TOOLS에 직접 등록한다.

사용법 (YAML spec):
  tools:
    - web_search
    - url_fetch
    - read_file
"""
from __future__ import annotations

import importlib
import json
import logging
from dataclasses import dataclass, field
from pathlib import Path
from typing import Any, Callable, Awaitable

logger = logging.getLogger(__name__)

_TOOLS_DIR = Path(__file__).parent / 'tools'


_CONFIG_PATH = Path(__file__).parent.parent.parent / 'data' / 'tool_config.json'


def _load_config() -> dict[str, str]:
    try:
        return json.loads(_CONFIG_PATH.read_text('utf-8')) if _CONFIG_PATH.exists() else {}
    except Exception:
        return {}


def save_tool_config(updates: dict[str, str]) -> None:
    cfg = _load_config()
    cfg.update({k: v for k, v in updates.items() if v})
    for k in [k for k, v in updates.items() if not v]:
        cfg.pop(k, None)
    _CONFIG_PATH.parent.mkdir(parents=True, exist_ok=True)
    _CONFIG_PATH.write_text(json.dumps(cfg, ensure_ascii=False, indent=2), 'utf-8')


def get_tool_config() -> dict[str, str]:
    return _load_config()


@dataclass
class ToolSpec:
    id: str
    name: str
    description: str
    category: str = 'general'      # general | research | file | code | design | integration
    is_async: bool = False
    enabled: bool = True
    params: list[str] = field(default_factory=list)  # context key들이 파라미터
    env_var: str = ''               # 필요한 환경변수/토큰 키 (예: FIGMA_TOKEN)


# ── 내장 도구 목록 ─────────────────────────────────────────────────────────────

_BUILTIN_TOOLS: dict[str, ToolSpec] = {
    # 'current_date', 'job_context' → jobs/tools/current_date.py, job_context.py 로 이동됨
    'web_search': ToolSpec(
        id='web_search',
        name='웹 검색',
        description='DuckDuckGo 또는 설정된 검색 엔진으로 웹 검색을 수행한다.',
        category='research',
        params=['topic', 'plan'],
    ),
    # 'url_fetch' → jobs/tools/url_fetch.py 로 이동됨
    'read_file': ToolSpec(
        id='read_file',
        name='파일 읽기',
        description='context의 file_path 키에 지정된 파일을 읽는다.',
        category='file',
        params=['file_path'],
    ),
    'list_files': ToolSpec(
        id='list_files',
        name='파일 목록',
        description='context의 dir_path 키에 지정된 디렉토리의 파일 목록을 가져온다.',
        category='file',
        params=['dir_path'],
    ),
    'run_shell': ToolSpec(
        id='run_shell',
        name='셸 명령 실행',
        description='context의 command 키에 지정된 명령을 실행한다. 주의: 보안 검토 필요.',
        category='code',
        enabled=False,  # 기본 비활성화
        params=['command'],
    ),
    'screenshot_url': ToolSpec(
        id='screenshot_url',
        name='URL 스크린샷',
        description='context의 url 키에 지정된 웹 페이지를 스크린샷으로 캡처한다. 비주얼 리서치·ui_designer 스텝에 유용.',
        category='research',
        params=['url'],
    ),
    'rss_feed': ToolSpec(
        id='rss_feed',
        name='RSS 피드 수집',
        description='context의 feed_url 키에 지정된 RSS/Atom 피드에서 최신 항목을 수집한다. 트렌드·뉴스 리서치에 유용.',
        category='research',
        params=['feed_url'],
    ),
    'write_file': ToolSpec(
        id='write_file',
        name='파일 저장',
        description='context의 file_path와 content 키를 사용해 파일을 저장한다. 산출물을 직접 파일로 출력할 때 사용.',
        category='file',
        params=['file_path', 'content'],
    ),
    'diff_files': ToolSpec(
        id='diff_files',
        name='파일 비교(diff)',
        description='context의 file_a, file_b 키에 지정된 두 파일의 diff를 반환한다. 코드 리뷰 시 변경사항 파악에 유용.',
        category='file',
        params=['file_a', 'file_b'],
    ),
    'run_python': ToolSpec(
        id='run_python',
        name='Python 코드 실행',
        description='context의 code 키에 담긴 Python 코드를 실행하고 결과를 반환한다. 데이터 분석·계산에 유용.',
        category='code',
        params=['code'],
    ),
    'figma_get_design': ToolSpec(
        id='figma_get_design',
        name='Figma 디자인 조회',
        description='context의 file_key와 node_id로 Figma 노드의 디자인 정보를 가져온다.',
        category='design',
        params=['file_key', 'node_id'],
        env_var='FIGMA_TOKEN',
    ),
    'figma_create_frame': ToolSpec(
        id='figma_create_frame',
        name='Figma 프레임 생성',
        description='[비활성] Figma REST API는 노드 직접 생성을 지원하지 않음. Figma Plugin API 필요.',
        category='design',
        enabled=False,
        params=['file_key', 'name', 'width', 'height'],
        env_var='FIGMA_TOKEN',
    ),
    'figma_update_node': ToolSpec(
        id='figma_update_node',
        name='Figma 노드 수정',
        description='[비활성] Figma REST API로 노드 속성 쓰기는 지원되지 않음. Figma Plugin API 필요.',
        category='design',
        enabled=False,
        params=['file_key', 'node_id', 'properties'],
        env_var='FIGMA_TOKEN',
    ),
    'slack_post': ToolSpec(
        id='slack_post',
        name='Slack 메시지 전송',
        description='context의 channel과 message를 Slack에 전송한다.',
        category='integration',
        params=['channel', 'message'],
        env_var='SLACK_BOT_TOKEN',
    ),
    'notion_write': ToolSpec(
        id='notion_write',
        name='Notion 페이지 작성',
        description='context의 page_id와 content를 Notion 페이지에 추가한다.',
        category='integration',
        params=['page_id', 'content'],
        env_var='NOTION_TOKEN',
    ),
    'pdf_generate': ToolSpec(
        id='pdf_generate',
        name='PDF 생성',
        description='context의 content(HTML 또는 Markdown)를 PDF 파일로 저장한다. context의 output_path에 저장 경로 지정.',
        category='file',
        params=['content', 'output_path'],
    ),
    'image_generate': ToolSpec(
        id='image_generate',
        name='이미지 생성',
        description='context의 prompt로 이미지를 생성한다. OPENAI_API_KEY 환경변수 필요. output_path에 PNG 저장.',
        category='file',
        params=['prompt', 'output_path'],
        env_var='OPENAI_API_KEY',
    ),
    'docx_generate': ToolSpec(
        id='docx_generate',
        name='Word 문서 생성',
        description='context의 content(Markdown)를 .docx 파일로 저장한다. output_path에 저장 경로 지정.',
        category='file',
        params=['content', 'output_path'],
    ),
    'pptx_generate': ToolSpec(
        id='pptx_generate',
        name='PowerPoint 생성',
        description='context의 content(마크다운 슬라이드 형식)를 .pptx 파일로 저장한다. ## 제목이 슬라이드 구분자. output_path에 저장 경로 지정.',
        category='file',
        params=['content', 'output_path'],
    ),
    'html_to_image': ToolSpec(
        id='html_to_image',
        name='HTML → 이미지',
        description='context의 html을 렌더링해 PNG로 저장한다. Playwright 설치 필요. output_path에 저장.',
        category='file',
        params=['html', 'output_path'],
    ),
    'webhook_call': ToolSpec(
        id='webhook_call',
        name='웹훅 / HTTP 호출',
        description='임의 URL로 HTTP 요청(GET/POST/PUT/DELETE)을 보낸다. url, method, body(JSON 문자열), headers(JSON) 지원.',
        category='integration',
        params=['url', 'method', 'body', 'headers'],
    ),
    'spreadsheet_read': ToolSpec(
        id='spreadsheet_read',
        name='스프레드시트 읽기',
        description='CSV 또는 Excel(.xlsx) 파일을 읽어 표 형태로 반환한다. file_path 필수. sheet(시트명), max_rows(기본 100) 선택.',
        category='file',
        params=['file_path', 'sheet', 'max_rows'],
    ),
    'email_send': ToolSpec(
        id='email_send',
        name='이메일 전송',
        description='SMTP로 이메일을 전송한다. to, subject, body 필수. html 선택. SMTP_USER·SMTP_PASS 환경변수 필요.',
        category='integration',
        params=['to', 'subject', 'body', 'html'],
        env_var='SMTP_PASS',
    ),
    'translate': ToolSpec(
        id='translate',
        name='텍스트 번역',
        description='text를 target_lang 언어로 번역한다. DEEPL_API_KEY 또는 GOOGLE_TRANSLATE_API_KEY 필요.',
        category='general',
        params=['text', 'target_lang', 'source_lang'],
        env_var='DEEPL_API_KEY',
    ),
    'airtable_query': ToolSpec(
        id='airtable_query',
        name='Airtable 조회/추가',
        description='Airtable 테이블을 조회(list)하거나 레코드를 추가(create)한다. base_id, table_name 필수. AIRTABLE_API_KEY 필요.',
        category='integration',
        params=['base_id', 'table_name', 'action', 'filter', 'fields'],
        env_var='AIRTABLE_API_KEY',
    ),
    'google_drive_upload': ToolSpec(
        id='google_drive_upload',
        name='Google Drive 파일 업로드',
        description='file_path의 파일을 Google Drive에 업로드한다. folder_id, file_name 선택. GOOGLE_SERVICE_ACCOUNT_JSON 필요.',
        category='integration',
        params=['file_path', 'folder_id', 'file_name'],
        env_var='GOOGLE_SERVICE_ACCOUNT_JSON',
    ),
    'calendar_create': ToolSpec(
        id='calendar_create',
        name='Google Calendar 일정 생성',
        description='Google Calendar에 일정을 생성한다. title, start, end(ISO 8601) 필수. GOOGLE_SERVICE_ACCOUNT_JSON 필요.',
        category='integration',
        params=['title', 'start', 'end', 'description', 'calendar_id'],
        env_var='GOOGLE_SERVICE_ACCOUNT_JSON',
    ),
}


def list_tools() -> list[dict[str, Any]]:
    """등록된 모든 도구 목록 반환 (플러그인 + legacy builtin 병합)."""
    from jobs.tools import load_plugin_tools
    plugins = load_plugin_tools()

    # 플러그인이 같은 id를 제공하면 legacy builtin 덮어쓰기
    merged: dict[str, ToolSpec] = dict(_BUILTIN_TOOLS)
    for tid, (spec, _fn) in plugins.items():
        merged[tid] = spec

    cfg = _load_config()
    import os
    return [
        {
            'id': t.id,
            'name': t.name,
            'description': t.description,
            'category': t.category,
            'enabled': t.enabled,
            'is_async': t.is_async,
            'params': t.params,
            'env_var': t.env_var,
            'token_set': bool(t.env_var and (os.environ.get(t.env_var) or cfg.get(t.env_var))),
        }
        for t in merged.values()
    ]


def execute_tool(tool_id: str, context: dict[str, str]) -> str:
    """동기 도구 실행 — plugin 우선, legacy builtin fallback."""
    # 플러그인 우선 — jobs/tools/<id>.py
    try:
        from jobs.tools import load_plugin_tools
        plugins = load_plugin_tools()
        if tool_id in plugins:
            _spec, exec_fn = plugins[tool_id]
            return exec_fn(context)
    except Exception as e:
        logger.debug('plugin 조회 실패 %s: %s', tool_id, e)

    # Legacy 내장 도구 (점진 이동 중)
    if tool_id == 'web_search':
        return _web_search(context)
    if tool_id == 'read_file':
        return _read_file(context)
    if tool_id == 'list_files':
        return _list_files(context)
    if tool_id == 'run_shell':
        spec = _BUILTIN_TOOLS.get('run_shell')
        if spec and not spec.enabled:
            return '[run_shell 비활성화됨]'
        return _run_shell(context)
    if tool_id == 'screenshot_url':
        return _screenshot_url(context)
    if tool_id == 'rss_feed':
        return _rss_feed(context)
    if tool_id == 'write_file':
        return _write_file(context)
    if tool_id == 'diff_files':
        return _diff_files(context)
    if tool_id == 'run_python':
        return _run_python(context)
    if tool_id == 'figma_get_design':
        return _figma_get_design(context)
    if tool_id == 'figma_create_frame':
        return _figma_create_frame(context)
    if tool_id == 'figma_update_node':
        return _figma_update_node(context)
    if tool_id == 'slack_post':
        return _slack_post(context)
    if tool_id == 'notion_write':
        return _notion_write(context)
    if tool_id == 'pdf_generate':
        return _pdf_generate(context)
    if tool_id == 'image_generate':
        return _image_generate(context)
    if tool_id == 'docx_generate':
        return _docx_generate(context)
    if tool_id == 'pptx_generate':
        return _pptx_generate(context)
    if tool_id == 'html_to_image':
        return _html_to_image(context)
    if tool_id == 'webhook_call':
        return _webhook_call(context)
    if tool_id == 'spreadsheet_read':
        return _spreadsheet_read(context)
    if tool_id == 'email_send':
        return _email_send(context)
    if tool_id == 'translate':
        return _translate(context)
    if tool_id == 'airtable_query':
        return _airtable_query(context)
    if tool_id == 'google_drive_upload':
        return _google_drive_upload(context)
    if tool_id == 'calendar_create':
        return _calendar_create(context)

    # 플러그인 도구
    if _TOOLS_DIR.exists():
        plugin_path = _TOOLS_DIR / f'{tool_id}.py'
        if plugin_path.exists():
            try:
                mod = importlib.import_module(f'jobs.tools.{tool_id}')
                if hasattr(mod, 'execute'):
                    return mod.execute(context)
            except Exception as e:
                logger.warning('플러그인 도구 실행 실패: %s — %s', tool_id, e)
                return f'[{tool_id} 실행 실패: {e}]'

    logger.warning('알 수 없는 도구: %s', tool_id)
    return ''


# ── 내장 도구 구현 ─────────────────────────────────────────────────────────────

def _web_search(context: dict[str, str]) -> str:
    import json as _json, re

    plan_text = context.get('plan', '')
    queries: list[str] = []
    try:
        m = re.search(r'\{[\s\S]*\}', plan_text)
        if m:
            plan_data = _json.loads(m.group())
            queries = plan_data.get('queries', [])
    except Exception:
        pass
    if not queries:
        queries = [context.get('topic', context.get('project', ''))]
    queries = [q for q in queries[:3] if q]

    brave_key = _resolve_token('BRAVE_SEARCH_API_KEY')
    if brave_key:
        return _web_search_brave(queries, brave_key)

    from harness.file_reader import web_search
    parts = [web_search(q, max_results=5) for q in queries]
    return '\n\n'.join(parts)


def _web_search_brave(queries: list[str], api_key: str) -> str:
    import urllib.request, urllib.parse, json as _json
    results = []
    for q in queries:
        encoded = urllib.parse.quote_plus(q)
        req = urllib.request.Request(
            f'https://api.search.brave.com/res/v1/web/search?q={encoded}&count=5',
            headers={'Accept': 'application/json', 'X-Subscription-Token': api_key},
        )
        try:
            with urllib.request.urlopen(req, timeout=10) as resp:
                data = _json.loads(resp.read())
            items = data.get('web', {}).get('results', [])
            parts = [f'## 검색: {q}']
            for item in items:
                title = item.get('title', '')
                url   = item.get('url', '')
                desc  = item.get('description', '')[:200]
                parts.append(f'- **{title}**\n  {url}\n  {desc}')
            results.append('\n'.join(parts))
        except Exception as e:
            results.append(f'[Brave 검색 실패: {q} — {e}]')
    return '\n\n'.join(results)


def _read_file(context: dict[str, str]) -> str:
    file_path = context.get('file_path', '')
    if not file_path:
        return ''
    try:
        p = Path(file_path)
        if not p.exists():
            return f'[파일 없음: {file_path}]'
        content = p.read_text(encoding='utf-8')
        # 너무 크면 앞부분만
        if len(content) > 20000:
            content = content[:20000] + f'\n\n[이후 {len(content) - 20000}자 생략]'
        return content
    except Exception as e:
        return f'[파일 읽기 실패: {e}]'


def _list_files(context: dict[str, str]) -> str:
    dir_path = context.get('dir_path', '.')
    try:
        p = Path(dir_path)
        if not p.is_dir():
            return f'[디렉토리 없음: {dir_path}]'
        files = [str(f.relative_to(p)) for f in sorted(p.rglob('*')) if f.is_file()][:100]
        return '\n'.join(files)
    except Exception as e:
        return f'[목록 실패: {e}]'


def _run_shell(context: dict[str, str]) -> str:
    import subprocess
    command = context.get('command', '')
    if not command:
        return ''
    try:
        result = subprocess.run(
            command, shell=True, capture_output=True, text=True, timeout=30,
        )
        output = result.stdout + result.stderr
        return output[:5000]
    except subprocess.TimeoutExpired:
        return '[명령 타임아웃 (30초)]'
    except Exception as e:
        return f'[명령 실행 실패: {e}]'


def _screenshot_url(context: dict[str, str]) -> str:
    import subprocess, os, time
    url = context.get('url', '')
    if not url:
        return '[screenshot_url: url 없음]'
    output_path = context.get('output_path', '')
    if not output_path:
        safe = ''.join(c if c.isalnum() else '_' for c in url)[:40]
        output_path = f'/tmp/screenshot_{safe}_{int(time.time())}.png'
    script = (
        'import os; from playwright.sync_api import sync_playwright; '
        'pw = sync_playwright().start(); b = pw.chromium.launch(); '
        'page = b.new_page(viewport={"width":1280,"height":800}); '
        'page.goto(os.environ["_SS_URL"], timeout=15000); '
        'page.screenshot(path=os.environ["_SS_OUT"], full_page=True); '
        'b.close(); pw.stop()'
    )
    env = {**os.environ, '_SS_URL': url, '_SS_OUT': output_path}
    try:
        Path(output_path).parent.mkdir(parents=True, exist_ok=True)
        result = subprocess.run(
            ['python3', '-c', script],
            capture_output=True, text=True, timeout=30, env=env,
        )
        if result.returncode == 0 and os.path.exists(output_path):
            size = Path(output_path).stat().st_size // 1024
            return f'[스크린샷 완료: {url}]\n저장 경로: {output_path} ({size}KB)'
        return f'[screenshot_url 실패: {result.stderr[:300]}]'
    except Exception as e:
        return f'[screenshot_url 실패: {e}]'


def _rss_feed(context: dict[str, str]) -> str:
    import urllib.request, xml.etree.ElementTree as ET
    feed_url = context.get('feed_url', '')
    if not feed_url:
        return '[rss_feed: feed_url 없음]'
    try:
        req = urllib.request.Request(feed_url, headers={'User-Agent': 'Mozilla/5.0'})
        with urllib.request.urlopen(req, timeout=10) as resp:
            raw = resp.read().decode('utf-8', errors='replace')
        root = ET.fromstring(raw)
        ns = {'atom': 'http://www.w3.org/2005/Atom'}
        items = []
        # RSS 2.0
        for item in root.findall('.//item')[:10]:
            title = item.findtext('title', '').strip()
            link  = item.findtext('link', '').strip()
            desc  = item.findtext('description', '').strip()[:200]
            items.append(f'- {title}\n  {link}\n  {desc}')
        # Atom
        if not items:
            for entry in root.findall('atom:entry', ns)[:10]:
                title = entry.findtext('atom:title', '', ns).strip()
                link_el = entry.find('atom:link', ns)
                link = link_el.get('href', '') if link_el is not None else ''
                items.append(f'- {title}\n  {link}')
        return f'[RSS: {feed_url}]\n\n' + '\n\n'.join(items) if items else '[RSS: 항목 없음]'
    except Exception as e:
        return f'[rss_feed 실패: {e}]'


def _write_file(context: dict[str, str]) -> str:
    file_path = context.get('file_path', '')
    content   = context.get('content', '')
    if not file_path:
        return '[write_file: file_path 없음]'
    try:
        p = Path(file_path)
        p.parent.mkdir(parents=True, exist_ok=True)
        p.write_text(content, encoding='utf-8')
        return f'[파일 저장 완료: {file_path} ({len(content)}자)]'
    except Exception as e:
        return f'[write_file 실패: {e}]'


def _diff_files(context: dict[str, str]) -> str:
    import difflib
    file_a = context.get('file_a', '')
    file_b = context.get('file_b', '')
    if not file_a or not file_b:
        return '[diff_files: file_a, file_b 모두 필요]'
    try:
        a_lines = Path(file_a).read_text(encoding='utf-8').splitlines(keepends=True)
        b_lines = Path(file_b).read_text(encoding='utf-8').splitlines(keepends=True)
        diff = list(difflib.unified_diff(a_lines, b_lines, fromfile=file_a, tofile=file_b))
        if not diff:
            return '[두 파일이 동일합니다]'
        result = ''.join(diff)
        return result[:8000] + (f'\n[이후 생략]' if len(result) > 8000 else '')
    except Exception as e:
        return f'[diff_files 실패: {e}]'


def _run_python(context: dict[str, str]) -> str:
    import subprocess, tempfile, os
    code = context.get('code', '')
    if not code:
        return '[run_python: code 없음]'
    try:
        with tempfile.NamedTemporaryFile(mode='w', suffix='.py', delete=False, encoding='utf-8') as f:
            f.write(code)
            tmp_path = f.name
        result = subprocess.run(
            ['python3', tmp_path],
            capture_output=True, text=True, timeout=30,
        )
        os.unlink(tmp_path)
        output = result.stdout + result.stderr
        return output[:5000] if output else '[출력 없음]'
    except subprocess.TimeoutExpired:
        return '[run_python 타임아웃 (30초)]'
    except Exception as e:
        return f'[run_python 실패: {e}]'


def _resolve_token(key: str) -> str:
    import os
    return os.environ.get(key) or _load_config().get(key, '')


def _figma_get_design(context: dict[str, str]) -> str:
    import urllib.request, json as _json
    token = _resolve_token('FIGMA_TOKEN')
    if not token:
        return '[figma_get_design: FIGMA_TOKEN 미설정 — 컴포넌트 라이브러리에서 토큰을 등록하세요]'
    file_key = context.get('file_key', '')
    node_id  = context.get('node_id', '')
    if not file_key:
        return '[figma_get_design: file_key 없음]'
    try:
        url = f'https://api.figma.com/v1/files/{file_key}'
        if node_id:
            url += f'/nodes?ids={node_id}'
        req = urllib.request.Request(url, headers={'X-Figma-Token': token})
        with urllib.request.urlopen(req, timeout=15) as resp:
            data = _json.loads(resp.read())
        return _json.dumps(data, ensure_ascii=False, indent=2)[:6000]
    except Exception as e:
        return f'[figma_get_design 실패: {e}]'


def _figma_create_frame(context: dict[str, str]) -> str:
    return '[figma_create_frame 비활성] Figma REST API는 노드 생성을 지원하지 않습니다. Figma Plugin API가 필요합니다.'


def _figma_update_node(context: dict[str, str]) -> str:
    return '[figma_update_node 비활성] Figma REST API로 노드 속성 쓰기는 지원되지 않습니다. Figma Plugin API가 필요합니다.'


def _slack_post(context: dict[str, str]) -> str:
    import urllib.request, json as _json
    token = _resolve_token('SLACK_BOT_TOKEN')
    if not token:
        return '[slack_post: SLACK_BOT_TOKEN 미설정 — 컴포넌트 라이브러리에서 토큰을 등록하세요]'
    channel = context.get('channel', '')
    message = context.get('message', '')
    if not channel or not message:
        return '[slack_post: channel, message 필요]'
    try:
        payload = _json.dumps({'channel': channel, 'text': message}).encode()
        req = urllib.request.Request(
            'https://slack.com/api/chat.postMessage',
            data=payload,
            headers={'Authorization': f'Bearer {token}', 'Content-Type': 'application/json'},
        )
        with urllib.request.urlopen(req, timeout=10) as resp:
            data = _json.loads(resp.read())
        if data.get('ok'):
            return f'[Slack 전송 완료: #{channel}]'
        return f'[slack_post 실패: {data.get("error", "unknown")}]'
    except Exception as e:
        return f'[slack_post 실패: {e}]'


def _pdf_generate(context: dict[str, str]) -> str:
    content     = context.get('content', '')
    output_path = context.get('output_path', '/tmp/output.pdf')
    if not content:
        return '[pdf_generate: content 없음]'
    try:
        # weasyprint 우선, 없으면 playwright fallback
        try:
            import weasyprint  # type: ignore
            import tempfile, os
            # Markdown → HTML 변환 시도
            try:
                import markdown as _md  # type: ignore
                html = _md.markdown(content, extensions=['tables', 'fenced_code'])
            except ImportError:
                html = f'<pre>{content}</pre>'
            html_full = f'<html><body style="font-family:sans-serif;padding:2em">{html}</body></html>'
            p = Path(output_path)
            p.parent.mkdir(parents=True, exist_ok=True)
            weasyprint.HTML(string=html_full).write_pdf(str(p))
            return f'[PDF 생성 완료: {output_path} ({p.stat().st_size // 1024}KB)]'
        except ImportError:
            pass
        # playwright fallback
        import subprocess, tempfile, os
        with tempfile.NamedTemporaryFile(mode='w', suffix='.html', delete=False, encoding='utf-8') as f:
            f.write(f'<html><body style="font-family:sans-serif;padding:2em"><pre>{content}</pre></body></html>')
            html_path = f.name
        p = Path(output_path)
        p.parent.mkdir(parents=True, exist_ok=True)
        result = subprocess.run(
            ['python3', '-c',
             f'from playwright.sync_api import sync_playwright; '
             f'pw = sync_playwright().start(); b = pw.chromium.launch(); '
             f'page = b.new_page(); page.goto("file://{html_path}"); '
             f'page.pdf(path="{output_path}"); b.close(); pw.stop()'],
            capture_output=True, text=True, timeout=30,
        )
        os.unlink(html_path)
        if result.returncode == 0:
            return f'[PDF 생성 완료: {output_path}]'
        return f'[pdf_generate 실패: {result.stderr[:300]}]'
    except Exception as e:
        return f'[pdf_generate 실패: {e}]'


def _image_generate(context: dict[str, str]) -> str:
    import urllib.request, json as _json, base64, os
    token = _resolve_token('OPENAI_API_KEY')
    if not token:
        return '[image_generate: OPENAI_API_KEY 미설정 — 컴포넌트 라이브러리에서 토큰을 등록하세요]'
    prompt      = context.get('prompt', '')
    output_path = context.get('output_path', '/tmp/output.png')
    if not prompt:
        return '[image_generate: prompt 없음]'
    try:
        payload = _json.dumps({
            'model': 'dall-e-3',
            'prompt': prompt,
            'n': 1,
            'size': '1024x1024',
            'response_format': 'b64_json',
        }).encode()
        req = urllib.request.Request(
            'https://api.openai.com/v1/images/generations',
            data=payload,
            headers={'Authorization': f'Bearer {token}', 'Content-Type': 'application/json'},
        )
        with urllib.request.urlopen(req, timeout=60) as resp:
            data = _json.loads(resp.read())
        b64 = data['data'][0]['b64_json']
        p = Path(output_path)
        p.parent.mkdir(parents=True, exist_ok=True)
        p.write_bytes(base64.b64decode(b64))
        return f'[이미지 생성 완료: {output_path} (1024×1024)]'
    except Exception as e:
        return f'[image_generate 실패: {e}]'


def _docx_generate(context: dict[str, str]) -> str:
    import re as _re
    content     = context.get('content', '')
    output_path = context.get('output_path', '/tmp/output.docx')
    if not content:
        return '[docx_generate: content 없음]'
    try:
        from docx import Document  # type: ignore
        from docx.shared import Pt, RGBColor  # type: ignore

        def _add_inline(para, text: str) -> None:
            parts = _re.split(r'(\*\*[^*\n]+\*\*|\*[^*\n]+\*|`[^`\n]+`)', text)
            for part in parts:
                if part.startswith('**') and part.endswith('**'):
                    run = para.add_run(part[2:-2])
                    run.bold = True
                elif part.startswith('*') and part.endswith('*'):
                    run = para.add_run(part[1:-1])
                    run.italic = True
                elif part.startswith('`') and part.endswith('`'):
                    run = para.add_run(part[1:-1])
                    run.font.name = 'Courier New'
                    run.font.size = Pt(10)
                elif part:
                    para.add_run(part)

        doc = Document()
        in_code = False
        code_buf: list[str] = []
        for line in content.splitlines():
            if line.startswith('```'):
                if in_code:
                    para = doc.add_paragraph(style='No Spacing')
                    run = para.add_run('\n'.join(code_buf))
                    run.font.name = 'Courier New'
                    run.font.size = Pt(10)
                    code_buf = []
                    in_code = False
                else:
                    in_code = True
                continue
            if in_code:
                code_buf.append(line)
                continue
            stripped = line.strip()
            if stripped.startswith('# '):
                doc.add_heading(stripped[2:], level=1)
            elif stripped.startswith('## '):
                doc.add_heading(stripped[3:], level=2)
            elif stripped.startswith('### '):
                doc.add_heading(stripped[4:], level=3)
            elif stripped.startswith('- ') or stripped.startswith('* '):
                para = doc.add_paragraph(style='List Bullet')
                _add_inline(para, stripped[2:])
            elif stripped:
                para = doc.add_paragraph()
                _add_inline(para, stripped)

        p = Path(output_path)
        p.parent.mkdir(parents=True, exist_ok=True)
        doc.save(str(p))
        return f'[Word 문서 생성 완료: {output_path}]'
    except ImportError:
        return '[docx_generate: python-docx 미설치 — pip install python-docx]'
    except Exception as e:
        return f'[docx_generate 실패: {e}]'


def _pptx_generate(context: dict[str, str]) -> str:
    content     = context.get('content', '')
    output_path = context.get('output_path', '/tmp/output.pptx')
    if not content:
        return '[pptx_generate: content 없음]'
    try:
        from pptx import Presentation  # type: ignore
        from pptx.util import Inches, Pt  # type: ignore

        prs = Presentation()
        slide_layout_title = prs.slide_layouts[0]  # 제목 슬라이드
        slide_layout_content = prs.slide_layouts[1]  # 제목+내용

        current_title = ''
        current_bullets: list[str] = []

        def _flush_slide() -> None:
            if not current_title and not current_bullets:
                return
            layout = slide_layout_title if not current_bullets else slide_layout_content
            slide = prs.slides.add_slide(layout)
            placeholders = slide.placeholders
            if placeholders:
                placeholders[0].text = current_title
            if current_bullets and len(placeholders) > 1:
                tf = placeholders[1].text_frame
                tf.clear()
                for i, bullet in enumerate(current_bullets):
                    p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
                    p.text = bullet

        for line in content.splitlines():
            stripped = line.strip()
            if stripped.startswith('## ') or stripped.startswith('# '):
                _flush_slide()
                current_title = stripped.lstrip('#').strip()
                current_bullets = []
            elif stripped.startswith('- ') or stripped.startswith('* '):
                current_bullets.append(stripped[2:])
            elif stripped:
                current_bullets.append(stripped)

        _flush_slide()

        p = Path(output_path)
        p.parent.mkdir(parents=True, exist_ok=True)
        prs.save(str(p))
        return f'[PowerPoint 생성 완료: {output_path}]'
    except ImportError:
        return '[pptx_generate: python-pptx 미설치 — pip install python-pptx]'
    except Exception as e:
        return f'[pptx_generate 실패: {e}]'


def _html_to_image(context: dict[str, str]) -> str:
    import subprocess, tempfile, os
    html        = context.get('html', '')
    output_path = context.get('output_path', '/tmp/output.png')
    if not html:
        return '[html_to_image: html 없음]'
    try:
        with tempfile.NamedTemporaryFile(mode='w', suffix='.html', delete=False, encoding='utf-8') as f:
            f.write(html)
            html_path = f.name
        p = Path(output_path)
        p.parent.mkdir(parents=True, exist_ok=True)
        result = subprocess.run(
            ['python3', '-c',
             f'from playwright.sync_api import sync_playwright; '
             f'pw = sync_playwright().start(); b = pw.chromium.launch(); '
             f'page = b.new_page(viewport={{"width":1280,"height":800}}); '
             f'page.goto("file://{html_path}"); '
             f'page.screenshot(path="{output_path}", full_page=True); '
             f'b.close(); pw.stop()'],
            capture_output=True, text=True, timeout=30,
        )
        os.unlink(html_path)
        if result.returncode == 0:
            size = Path(output_path).stat().st_size // 1024
            return f'[HTML → 이미지 완료: {output_path} ({size}KB)]'
        return f'[html_to_image 실패: {result.stderr[:300]}]'
    except Exception as e:
        return f'[html_to_image 실패: {e}]'


def _webhook_call(context: dict[str, str]) -> str:
    import urllib.request, json as _json
    url     = context.get('url', '')
    method  = context.get('method', 'POST').upper()
    body_raw = context.get('body', '')
    headers_raw = context.get('headers', '{}')
    if not url:
        return '[webhook_call: url 없음]'
    try:
        extra_headers = _json.loads(headers_raw) if headers_raw else {}
    except Exception:
        extra_headers = {}
    body_bytes = None
    if body_raw:
        if 'Content-Type' not in extra_headers:
            extra_headers['Content-Type'] = 'application/json'
        body_bytes = body_raw.encode() if isinstance(body_raw, str) else body_raw
    try:
        req = urllib.request.Request(url, data=body_bytes, headers=extra_headers, method=method)
        with urllib.request.urlopen(req, timeout=15) as resp:
            status = resp.status
            data = resp.read().decode('utf-8', errors='replace')[:2000]
        return f'[webhook_call 완료: {method} {url} — HTTP {status}]\n{data}'
    except Exception as e:
        return f'[webhook_call 실패: {e}]'


def _spreadsheet_read(context: dict[str, str]) -> str:
    import csv, io
    file_path = context.get('file_path', '')
    sheet     = context.get('sheet', '')
    try:
        max_rows = int(context.get('max_rows', '100'))
    except ValueError:
        max_rows = 100
    if not file_path:
        return '[spreadsheet_read: file_path 없음]'
    p = Path(file_path)
    if not p.exists():
        return f'[spreadsheet_read: 파일 없음: {file_path}]'
    ext = p.suffix.lower()
    if ext == '.csv':
        try:
            content = p.read_text(encoding='utf-8-sig')
            reader = csv.DictReader(io.StringIO(content))
            rows = list(reader)[:max_rows]
            if not rows:
                return '[spreadsheet_read: 데이터 없음]'
            cols = list(rows[0].keys())
            lines = ['\t'.join(cols)] + ['\t'.join(str(r.get(c, '')) for c in cols) for r in rows]
            return f'[CSV: {file_path} — {len(rows)}행]\n' + '\n'.join(lines)
        except Exception as e:
            return f'[spreadsheet_read CSV 실패: {e}]'
    elif ext in ('.xlsx', '.xls'):
        try:
            import openpyxl  # type: ignore
            wb = openpyxl.load_workbook(file_path, read_only=True, data_only=True)
            ws = wb[sheet] if sheet and sheet in wb.sheetnames else wb.active
            rows = []
            for row in ws.iter_rows(values_only=True):
                if len(rows) >= max_rows:
                    break
                rows.append([str(c) if c is not None else '' for c in row])
            if not rows:
                return '[spreadsheet_read: 데이터 없음]'
            lines = ['\t'.join(rows[0])] + ['\t'.join(r) for r in rows[1:]]
            return f'[Excel: {file_path} — {len(rows) - 1}행]\n' + '\n'.join(lines)
        except ImportError:
            return '[spreadsheet_read: openpyxl 미설치 — pip install openpyxl]'
        except Exception as e:
            return f'[spreadsheet_read Excel 실패: {e}]'
    return f'[spreadsheet_read: 지원하지 않는 형식 {ext} — .csv/.xlsx만 지원]'


def _email_send(context: dict[str, str]) -> str:
    import smtplib, os
    from email.mime.text import MIMEText
    from email.mime.multipart import MIMEMultipart
    smtp_host = _resolve_token('SMTP_HOST') or os.environ.get('SMTP_HOST', 'smtp.gmail.com')
    smtp_port_raw = _resolve_token('SMTP_PORT') or os.environ.get('SMTP_PORT', '587')
    smtp_user = _resolve_token('SMTP_USER') or os.environ.get('SMTP_USER', '')
    smtp_pass = _resolve_token('SMTP_PASS') or os.environ.get('SMTP_PASS', '')
    to      = context.get('to', '')
    subject = context.get('subject', '(제목 없음)')
    body    = context.get('body', '')
    html    = context.get('html', '')
    if not to:
        return '[email_send: to 없음]'
    if not smtp_user or not smtp_pass:
        return '[email_send: SMTP_USER, SMTP_PASS 미설정 — 컴포넌트 라이브러리에서 토큰을 등록하세요]'
    try:
        msg = MIMEMultipart('alternative')
        msg['Subject'] = subject
        msg['From'] = smtp_user
        msg['To'] = to
        if body:
            msg.attach(MIMEText(body, 'plain', 'utf-8'))
        if html:
            msg.attach(MIMEText(html, 'html', 'utf-8'))
        with smtplib.SMTP(smtp_host, int(smtp_port_raw)) as server:
            server.starttls()
            server.login(smtp_user, smtp_pass)
            server.sendmail(smtp_user, [t.strip() for t in to.split(',')], msg.as_string())
        return f'[email_send 완료: {to} — "{subject}"]'
    except Exception as e:
        return f'[email_send 실패: {e}]'


def _translate(context: dict[str, str]) -> str:
    import urllib.request, urllib.parse, json as _json
    text        = context.get('text', '')
    target_lang = context.get('target_lang', 'ko')
    source_lang = context.get('source_lang', '')
    if not text:
        return '[translate: text 없음]'
    deepl_key = _resolve_token('DEEPL_API_KEY')
    if deepl_key:
        try:
            params: dict = {'text': text, 'target_lang': target_lang.upper()}
            if source_lang:
                params['source_lang'] = source_lang.upper()
            payload = urllib.parse.urlencode(params).encode()
            base = 'api-free.deepl.com' if deepl_key.endswith(':fx') else 'api.deepl.com'
            req = urllib.request.Request(
                f'https://{base}/v2/translate', data=payload,
                headers={'Authorization': f'DeepL-Auth-Key {deepl_key}',
                         'Content-Type': 'application/x-www-form-urlencoded'},
            )
            with urllib.request.urlopen(req, timeout=15) as resp:
                data = _json.loads(resp.read())
            result = data['translations'][0]['text']
            detected = data['translations'][0].get('detected_source_language', '')
            return f'[번역 완료 ({detected}→{target_lang.upper()})]\n{result}'
        except Exception as e:
            return f'[translate DeepL 실패: {e}]'
    google_key = _resolve_token('GOOGLE_TRANSLATE_API_KEY')
    if google_key:
        try:
            payload = _json.dumps({
                'q': text, 'target': target_lang,
                'source': source_lang or None, 'format': 'text',
            }).encode()
            req = urllib.request.Request(
                f'https://translation.googleapis.com/language/translate/v2?key={google_key}',
                data=payload, headers={'Content-Type': 'application/json'},
            )
            with urllib.request.urlopen(req, timeout=15) as resp:
                data = _json.loads(resp.read())
            result = data['data']['translations'][0]['translatedText']
            return f'[번역 완료 →{target_lang}]\n{result}'
        except Exception as e:
            return f'[translate Google 실패: {e}]'
    return '[translate: DEEPL_API_KEY 또는 GOOGLE_TRANSLATE_API_KEY 미설정]'


def _airtable_query(context: dict[str, str]) -> str:
    import urllib.request, urllib.parse, json as _json
    token = _resolve_token('AIRTABLE_API_KEY')
    if not token:
        return '[airtable_query: AIRTABLE_API_KEY 미설정 — 컴포넌트 라이브러리에서 토큰을 등록하세요]'
    base_id    = context.get('base_id', '')
    table_name = context.get('table_name', '')
    action     = context.get('action', 'list')
    filter_formula = context.get('filter', '')
    fields_raw = context.get('fields', '{}')
    if not base_id or not table_name:
        return '[airtable_query: base_id, table_name 필요]'
    headers = {'Authorization': f'Bearer {token}', 'Content-Type': 'application/json'}
    base_url = f'https://api.airtable.com/v0/{base_id}/{urllib.parse.quote(table_name)}'
    if action == 'create':
        try:
            fields = _json.loads(fields_raw) if fields_raw else {}
            payload = _json.dumps({'fields': fields}).encode()
            req = urllib.request.Request(base_url, data=payload, headers=headers)
            with urllib.request.urlopen(req, timeout=10) as resp:
                data = _json.loads(resp.read())
            return f'[Airtable 레코드 생성: {data.get("id", "")}]\n{_json.dumps(data.get("fields", {}), ensure_ascii=False)[:500]}'
        except Exception as e:
            return f'[airtable_query create 실패: {e}]'
    else:
        params: dict = {'maxRecords': '50'}
        if filter_formula:
            params['filterByFormula'] = filter_formula
        url = base_url + '?' + urllib.parse.urlencode(params)
        try:
            req = urllib.request.Request(url, headers=headers)
            with urllib.request.urlopen(req, timeout=10) as resp:
                data = _json.loads(resp.read())
            records = data.get('records', [])
            if not records:
                return '[Airtable: 레코드 없음]'
            lines = [f'[Airtable: {table_name} — {len(records)}개]']
            for rec in records:
                flds = rec.get('fields', {})
                preview = ', '.join(f'{k}={v}' for k, v in list(flds.items())[:5])
                lines.append(f'- {rec["id"]}: {preview}')
            return '\n'.join(lines)
        except Exception as e:
            return f'[airtable_query list 실패: {e}]'


def _google_drive_upload(context: dict[str, str]) -> str:
    import json as _json
    sa_json_str = _resolve_token('GOOGLE_SERVICE_ACCOUNT_JSON')
    if not sa_json_str:
        return '[google_drive_upload: GOOGLE_SERVICE_ACCOUNT_JSON 미설정 — 컴포넌트 라이브러리에서 등록하세요]'
    file_path = context.get('file_path', '')
    folder_id = context.get('folder_id', '')
    file_name = context.get('file_name', '')
    if not file_path:
        return '[google_drive_upload: file_path 없음]'
    p = Path(file_path)
    if not p.exists():
        return f'[google_drive_upload: 파일 없음: {file_path}]'
    if not file_name:
        file_name = p.name
    try:
        import google.oauth2.service_account as _sa  # type: ignore
        import googleapiclient.discovery as _discovery  # type: ignore
        from googleapiclient.http import MediaFileUpload  # type: ignore
        creds = _sa.Credentials.from_service_account_info(
            _json.loads(sa_json_str),
            scopes=['https://www.googleapis.com/auth/drive.file'],
        )
        service = _discovery.build('drive', 'v3', credentials=creds)
        metadata: dict = {'name': file_name}
        if folder_id:
            metadata['parents'] = [folder_id]
        result = service.files().create(
            body=metadata, media_body=MediaFileUpload(file_path),
            fields='id,name,webViewLink',
        ).execute()
        return (f'[Google Drive 업로드 완료: {result.get("name")}]\n'
                f'링크: {result.get("webViewLink", "")}\nID: {result.get("id", "")}')
    except ImportError:
        return '[google_drive_upload: google-api-python-client 미설치 — pip install google-api-python-client google-auth]'
    except Exception as e:
        return f'[google_drive_upload 실패: {e}]'


def _calendar_create(context: dict[str, str]) -> str:
    import json as _json
    sa_json_str = _resolve_token('GOOGLE_SERVICE_ACCOUNT_JSON')
    if not sa_json_str:
        return '[calendar_create: GOOGLE_SERVICE_ACCOUNT_JSON 미설정 — 컴포넌트 라이브러리에서 등록하세요]'
    title       = context.get('title', '')
    start       = context.get('start', '')
    end         = context.get('end', '')
    description = context.get('description', '')
    calendar_id = context.get('calendar_id', 'primary')
    if not title or not start or not end:
        return '[calendar_create: title, start, end(ISO 8601) 필요]'
    try:
        import google.oauth2.service_account as _sa  # type: ignore
        import googleapiclient.discovery as _discovery  # type: ignore
        creds = _sa.Credentials.from_service_account_info(
            _json.loads(sa_json_str),
            scopes=['https://www.googleapis.com/auth/calendar'],
        )
        service = _discovery.build('calendar', 'v3', credentials=creds)
        event = {
            'summary': title,
            'description': description,
            'start': {'dateTime': start, 'timeZone': 'Asia/Seoul'},
            'end':   {'dateTime': end,   'timeZone': 'Asia/Seoul'},
        }
        result = service.events().insert(calendarId=calendar_id, body=event).execute()
        return (f'[일정 생성 완료: {title}]\n'
                f'링크: {result.get("htmlLink", "")}\nID: {result.get("id", "")}')
    except ImportError:
        return '[calendar_create: google-api-python-client 미설치 — pip install google-api-python-client google-auth]'
    except Exception as e:
        return f'[calendar_create 실패: {e}]'


def _notion_write(context: dict[str, str]) -> str:
    import urllib.request, json as _json
    token = _resolve_token('NOTION_TOKEN')
    if not token:
        return '[notion_write: NOTION_TOKEN 미설정 — 컴포넌트 라이브러리에서 토큰을 등록하세요]'
    page_id = context.get('page_id', '')
    content = context.get('content', '')
    if not page_id or not content:
        return '[notion_write: page_id, content 필요]'

    def _rich_text(text: str) -> list:
        return [{'type': 'text', 'text': {'content': text[:2000]}}]

    blocks: list[dict] = []
    code_buf: list[str] = []
    in_code = False
    for line in content.splitlines():
        if line.startswith('```'):
            if in_code:
                blocks.append({'object': 'block', 'type': 'code',
                                'code': {'rich_text': _rich_text('\n'.join(code_buf)), 'language': 'plain text'}})
                code_buf = []
                in_code = False
            else:
                in_code = True
            continue
        if in_code:
            code_buf.append(line)
            continue
        stripped = line.strip()
        if stripped.startswith('# '):
            blocks.append({'object': 'block', 'type': 'heading_1',
                           'heading_1': {'rich_text': _rich_text(stripped[2:])}})
        elif stripped.startswith('## '):
            blocks.append({'object': 'block', 'type': 'heading_2',
                           'heading_2': {'rich_text': _rich_text(stripped[3:])}})
        elif stripped.startswith('### '):
            blocks.append({'object': 'block', 'type': 'heading_3',
                           'heading_3': {'rich_text': _rich_text(stripped[4:])}})
        elif stripped.startswith('- ') or stripped.startswith('* '):
            blocks.append({'object': 'block', 'type': 'bulleted_list_item',
                           'bulleted_list_item': {'rich_text': _rich_text(stripped[2:])}})
        elif stripped.startswith('> '):
            blocks.append({'object': 'block', 'type': 'quote',
                           'quote': {'rich_text': _rich_text(stripped[2:])}})
        elif stripped:
            blocks.append({'object': 'block', 'type': 'paragraph',
                           'paragraph': {'rich_text': _rich_text(stripped)}})

    headers = {'Authorization': f'Bearer {token}', 'Content-Type': 'application/json',
                'Notion-Version': '2022-06-28'}
    total = 0
    try:
        for i in range(0, len(blocks), 100):
            payload = _json.dumps({'children': blocks[i:i + 100]}).encode()
            req = urllib.request.Request(
                f'https://api.notion.com/v1/blocks/{page_id}/children',
                data=payload, headers=headers,
            )
            with urllib.request.urlopen(req, timeout=10):
                pass
            total += len(blocks[i:i + 100])
        return f'[Notion 작성 완료: {page_id} — {total}개 블록]'
    except Exception as e:
        return f'[notion_write 실패: {e}]'
